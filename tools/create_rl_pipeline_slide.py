"""전국 MCI 강화학습 연구 파이프라인 1장 슬라이드를 생성한다."""

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/ryu/MCI_UAV")
OUT = ROOT / "docs" / "presentation"
TEMPLATE = Path(
    "/home/ryu/.codex/attachments/d37a2c0a-0977-4735-8229-eb9cdf807a64/"
    "codex-clipboard-0fb12bc6-9953-4dba-b6dc-a391faaa5211.png"
)
AGENT = Path(
    "/home/ryu/.codex/attachments/b6c175ed-88dc-4dce-bbef-b5776154ed51/"
    "codex-clipboard-c7a2c758-4acf-4b91-9391-e726654eea12.png"
)
PPTX_OUT = OUT / "MCI_RL_연구_파이프라인_1슬라이드.pptx"
PNG_OUT = OUT / "MCI_RL_연구_파이프라인_미리보기.png"

SW, SH = 13.333333, 7.5
PX_W, PX_H = 1600, 900
SCALE = 120.0

FONT = "맑은 고딕"
PIL_FONT = "/home/ryu/.fonts/NanumGothic-Regular.ttf"

NAVY = "173A55"
BLUE = "0A83BF"
CYAN = "16A4D8"
TEAL = "17A398"
ORANGE = "F09A3E"
GRAY = "71808E"
MID = "CDD7DF"
LIGHT = "F5F8FA"
WHITE = "FFFFFF"
BLACK = "17212B"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_text(slide, x, y, w, h, text, size=14, color=BLACK, bold=False,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return box


def add_round_rect(slide, x, y, w, h, fill=WHITE, line=MID, radius=True,
                   line_width=1.0):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_line(slide, x1, y1, x2, y2, color=CYAN, width=2.4, dash=None):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash is not None:
        line.line.dash_style = dash
    return line


def add_arrow(slide, x, y, w=0.30, h=0.17, color=CYAN):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def stage_label(slide, x, w, number, title):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(1.58), Inches(0.27), Inches(0.27)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(BLUE)
    circle.line.fill.background()
    add_text(slide, x, 1.58, 0.27, 0.27, str(number), 10, WHITE, True)
    add_text(slide, x + 0.34, 1.54, w - 0.34, 0.34, title, 15, NAVY, True,
             PP_ALIGN.LEFT)


def icon_scenario(slide, cx, cy):
    # 문서 3장과 위치 핀으로 다지역 시나리오 묶음을 표현한다.
    for dx, dy, fill in [(-0.22, 0.05, "EAF4FA"), (-0.10, -0.01, "DDF0F8"), (0.02, -0.07, WHITE)]:
        doc = add_round_rect(slide, cx + dx, cy + dy, 0.46, 0.55, fill, BLUE, True, 1.0)
        doc.adjustments[0] = 0.05
    for i in range(3):
        add_line(slide, cx + 0.10, cy + 0.06 + i * 0.10, cx + 0.34,
                 cy + 0.06 + i * 0.10, GRAY, 1.1)
    pin = slide.shapes.add_shape(
        MSO_SHAPE.TEAR, Inches(cx - 0.38), Inches(cy - 0.18), Inches(0.34), Inches(0.34)
    )
    pin.rotation = 225
    pin.fill.solid()
    pin.fill.fore_color.rgb = rgb(ORANGE)
    pin.line.fill.background()


def icon_simulator(slide, cx, cy):
    screen = add_round_rect(slide, cx - 0.36, cy - 0.20, 0.72, 0.47, WHITE, BLUE, True, 1.4)
    screen.adjustments[0] = 0.08
    add_line(slide, cx - 0.25, cy + 0.02, cx - 0.10, cy + 0.02, TEAL, 1.7)
    add_line(slide, cx - 0.10, cy + 0.02, cx - 0.03, cy - 0.09, TEAL, 1.7)
    add_line(slide, cx - 0.03, cy - 0.09, cx + 0.07, cy + 0.13, TEAL, 1.7)
    add_line(slide, cx + 0.07, cy + 0.13, cx + 0.16, cy - 0.02, TEAL, 1.7)
    add_line(slide, cx + 0.16, cy - 0.02, cx + 0.27, cy - 0.02, TEAL, 1.7)
    add_line(slide, cx, cy + 0.29, cx, cy + 0.37, GRAY, 1.4)
    add_line(slide, cx - 0.18, cy + 0.37, cx + 0.18, cy + 0.37, GRAY, 1.4)


def icon_checklist(slide, x, y):
    page = add_round_rect(slide, x, y, 0.45, 0.55, WHITE, GRAY, True, 1.1)
    page.adjustments[0] = 0.05
    for i in range(3):
        yy = y + 0.12 + i * 0.14
        add_line(slide, x + 0.09, yy + 0.02, x + 0.14, yy + 0.07, TEAL, 1.2)
        add_line(slide, x + 0.14, yy + 0.07, x + 0.21, yy - 0.01, TEAL, 1.2)
        add_line(slide, x + 0.25, yy + 0.03, x + 0.38, yy + 0.03, GRAY, 1.0)


def icon_holdout(slide, cx, cy):
    for d, color, width in [(0.56, BLUE, 1.3), (0.38, CYAN, 1.1), (0.20, ORANGE, 1.0)]:
        c = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d)
        )
        c.fill.background()
        c.line.color.rgb = rgb(color)
        c.line.width = Pt(width)
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - 0.045), Inches(cy - 0.045), Inches(0.09), Inches(0.09)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = rgb(ORANGE)
    dot.line.fill.background()


def icon_log(slide, x, y):
    page = add_round_rect(slide, x, y, 0.42, 0.49, WHITE, BLUE, True, 1.0)
    page.adjustments[0] = 0.05
    for i, w in enumerate([0.25, 0.20, 0.27]):
        add_line(slide, x + 0.08, y + 0.13 + i * 0.11, x + 0.08 + w,
                 y + 0.13 + i * 0.11, GRAY, 1.0)


def icon_tree(slide, cx, cy, color=TEAL):
    nodes = [
        (cx, cy - 0.19),
        (cx - 0.23, cy + 0.02), (cx + 0.23, cy + 0.02),
        (cx - 0.34, cy + 0.25), (cx - 0.12, cy + 0.25),
        (cx + 0.12, cy + 0.25), (cx + 0.34, cy + 0.25),
    ]
    for a, b in [(0, 1), (0, 2), (1, 3), (1, 4), (2, 5), (2, 6)]:
        add_line(slide, nodes[a][0], nodes[a][1], nodes[b][0], nodes[b][1], color, 1.4)
    for i, (x, y) in enumerate(nodes):
        s = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x - 0.055), Inches(y - 0.055), Inches(0.11), Inches(0.11)
        )
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(ORANGE if i == 0 else color)
        s.line.fill.background()


def icon_compare(slide, x, y):
    colors = [GRAY, BLUE, TEAL]
    widths = [0.57, 0.73, 0.64]
    for i, (color, width) in enumerate(zip(colors, widths)):
        yy = y + i * 0.25
        bar = add_round_rect(slide, x, yy, width, 0.13, color, color, True, 0.0)
        bar.adjustments[0] = 0.25


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.shapes.add_picture(str(TEMPLATE), 0, 0, width=Inches(SW), height=Inches(SH))
    panel = add_round_rect(slide, 0.44, 1.31, 12.56, 5.37, "F5F7F9", "CFD6DC", True, 1.0)
    panel.adjustments[0] = 0.025

    # 단계 제목
    stage_label(slide, 0.70, 2.03, 1, "시나리오·환경")
    stage_label(slide, 3.02, 2.75, 2, "정책 생성")
    stage_label(slide, 6.23, 1.90, 3, "일반화 검증")
    stage_label(slide, 8.43, 2.27, 4, "정책 증류")
    stage_label(slide, 11.04, 1.70, 5, "최종 비교")

    # 1. 시나리오·환경
    env_card = add_round_rect(slide, 0.67, 2.03, 2.10, 3.74, WHITE, MID, True, 1.0)
    env_card.adjustments[0] = 0.04
    icon_scenario(slide, 1.30, 2.37)
    add_text(slide, 0.88, 2.78, 1.68, 0.28, "전국 시나리오 생성", 15, NAVY, True)
    add_text(slide, 0.86, 3.07, 1.72, 0.36, "시군구 · 병원 · AMB/UAV", 10.5, GRAY)
    add_line(slide, 1.72, 3.52, 1.72, 3.81, CYAN, 2.0)
    down = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(1.59), Inches(3.69), Inches(0.26), Inches(0.29)
    )
    down.fill.solid(); down.fill.fore_color.rgb = rgb(CYAN); down.line.fill.background()
    icon_simulator(slide, 1.72, 4.30)
    add_text(slide, 0.87, 4.78, 1.70, 0.28, "사건기반 시뮬레이션", 15, NAVY, True)
    add_text(slide, 0.84, 5.09, 1.76, 0.42, "동일 환경에서 정책 실행\n환자 이송–병원 치료 반영", 10.5, GRAY)

    # 분기 화살표
    add_line(slide, 2.77, 3.88, 2.92, 3.88, CYAN, 2.8)
    add_line(slide, 2.92, 2.70, 2.92, 4.82, CYAN, 2.8)
    add_line(slide, 2.92, 2.70, 3.08, 2.70, CYAN, 2.8)
    add_line(slide, 2.92, 4.82, 3.08, 4.82, CYAN, 2.8)
    add_arrow(slide, 2.99, 2.61, 0.24, 0.18, CYAN)
    add_arrow(slide, 2.99, 4.73, 0.24, 0.18, CYAN)

    # 2. 휴리스틱과 RL 정책 생성
    heur = add_round_rect(slide, 3.14, 2.08, 2.65, 1.26, "F7FAFC", "BFCBD4", True, 1.0)
    heur.adjustments[0] = 0.06
    icon_checklist(slide, 3.36, 2.41)
    add_text(slide, 3.96, 2.23, 1.55, 0.28, "64개 휴리스틱", 15, NAVY, True, PP_ALIGN.LEFT)
    add_text(slide, 3.96, 2.54, 1.55, 0.44, "규칙별 기준 성능 생성\n→ 비교군 확보", 10.5, GRAY, False, PP_ALIGN.LEFT)

    rl = add_round_rect(slide, 3.14, 3.55, 2.65, 2.22, "EEF8FC", "7EC4E2", True, 1.4)
    rl.adjustments[0] = 0.05
    add_text(slide, 3.34, 3.70, 1.15, 0.30, "전국 RL 학습", 15.5, NAVY, True, PP_ALIGN.LEFT)
    add_text(slide, 3.34, 4.03, 1.10, 0.58, "시군구 다지역\nMaskable PPO", 10.5, GRAY, False, PP_ALIGN.LEFT)
    slide.shapes.add_picture(str(AGENT), Inches(4.48), Inches(3.77), width=Inches(1.07), height=Inches(1.09))
    add_round_rect(slide, 3.35, 5.05, 2.24, 0.43, "DDF2FA", "DDF2FA", True, 0.0)
    add_text(slide, 3.43, 5.10, 2.08, 0.31, "상태·보상  ↔  마스킹된 행동", 10, BLUE, True)

    # 두 정책을 Holdout으로 합류
    add_line(slide, 5.79, 2.70, 6.03, 2.70, CYAN, 2.3)
    add_line(slide, 5.79, 4.66, 6.03, 4.66, CYAN, 2.3)
    add_line(slide, 6.03, 2.70, 6.03, 4.66, CYAN, 2.3)
    add_line(slide, 6.03, 3.68, 6.25, 3.68, CYAN, 2.8)
    add_arrow(slide, 6.14, 3.59, 0.25, 0.18, CYAN)

    # 3. 학습에 쓰지 않은 좌표에서 쌍비교
    hold = add_round_rect(slide, 6.36, 2.22, 1.76, 3.35, WHITE, MID, True, 1.0)
    hold.adjustments[0] = 0.05
    icon_holdout(slide, 7.24, 2.85)
    add_text(slide, 6.53, 3.24, 1.42, 0.32, "Holdout 평가", 16, NAVY, True)
    add_text(slide, 6.53, 3.60, 1.42, 0.46, "학습하지 않은\n지역·좌표", 11, GRAY)
    add_round_rect(slide, 6.58, 4.22, 1.31, 0.43, "EEF3F6", "EEF3F6", True, 0.0)
    add_text(slide, 6.66, 4.27, 1.15, 0.31, "휴리스틱  ↔  RL", 10.5, BLUE, True)
    add_text(slide, 6.53, 4.78, 1.42, 0.46, "지역별 성능 지도\n동일 seed 통계 비교", 10.2, GRAY)

    add_line(slide, 8.12, 3.68, 8.42, 3.68, CYAN, 2.8)
    add_arrow(slide, 8.31, 3.59, 0.25, 0.18, CYAN)

    # 4. 로그→VIPER→Tree 재실행
    dist = add_round_rect(slide, 8.53, 2.05, 2.18, 3.72, WHITE, MID, True, 1.0)
    dist.adjustments[0] = 0.04
    icon_log(slide, 8.78, 2.36)
    add_text(slide, 9.32, 2.27, 1.16, 0.32, "RL 행동 로그", 13.5, NAVY, True, PP_ALIGN.LEFT)
    add_text(slide, 9.32, 2.58, 1.16, 0.24, "상태–행동 수집", 9.8, GRAY, False, PP_ALIGN.LEFT)
    add_line(slide, 9.62, 2.98, 9.62, 3.21, CYAN, 1.8)
    d1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.52), Inches(3.11), Inches(0.20), Inches(0.22))
    d1.fill.solid(); d1.fill.fore_color.rgb = rgb(CYAN); d1.line.fill.background()
    icon_tree(slide, 9.62, 3.66, TEAL)
    add_text(slide, 8.78, 4.02, 1.68, 0.28, "VIPER 정책 증류", 14.5, NAVY, True)
    add_text(slide, 8.78, 4.32, 1.68, 0.29, "RL → 의사결정트리", 10.2, GRAY)
    add_line(slide, 9.62, 4.68, 9.62, 4.91, CYAN, 1.8)
    d2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.52), Inches(4.81), Inches(0.20), Inches(0.22))
    d2.fill.solid(); d2.fill.fore_color.rgb = rgb(CYAN); d2.line.fill.background()
    add_round_rect(slide, 8.81, 5.08, 1.62, 0.46, "E7F6F3", "B5DDD7", True, 0.8)
    add_text(slide, 8.91, 5.13, 1.42, 0.33, "Tree 재시뮬레이션", 10.8, TEAL, True)

    add_line(slide, 10.71, 3.68, 11.02, 3.68, CYAN, 2.8)
    add_arrow(slide, 10.91, 3.59, 0.25, 0.18, CYAN)

    # 5. 세 정책을 같은 축에서 최종 평가
    final = add_round_rect(slide, 11.12, 2.11, 1.68, 3.56, WHITE, MID, True, 1.0)
    final.adjustments[0] = 0.05
    icon_compare(slide, 11.61, 2.48)
    add_text(slide, 11.31, 3.13, 1.30, 0.30, "3자 최종 비교", 15, NAVY, True)
    items = [("휴리스틱", "EEF1F3", GRAY), ("RL 정책", "E5F3FA", BLUE), ("VIPER Tree", "E4F5F2", TEAL)]
    for i, (label, fill, color) in enumerate(items):
        yy = 3.60 + i * 0.51
        add_round_rect(slide, 11.35, yy, 1.21, 0.37, fill, fill, True, 0.0)
        add_text(slide, 11.42, yy + 0.03, 1.07, 0.29, label, 10.5, color, True)
    add_text(slide, 11.31, 5.15, 1.30, 0.34, "생존 성과 · 일반화\n· 해석성", 9.8, GRAY)

    # 한 문장으로 연구 방향을 고정한다.
    banner = add_round_rect(slide, 0.84, 6.02, 11.78, 0.43, "E8F3F8", "C7E2EF", True, 0.7)
    banner.adjustments[0] = 0.18
    add_text(slide, 1.03, 6.07, 11.40, 0.31,
             "전국 시나리오에서 학습 → 미학습 지역에서 검증 → 해석 가능한 정책으로 환류",
             13, NAVY, True)

    # 이미지 기반 기존 페이지 번호 위에 새 번호를 덮지 않고 템플릿을 그대로 유지한다.
    prs.save(PPTX_OUT)


def pil_font(size):
    return ImageFont.truetype(PIL_FONT, size=size)


def draw_center(draw, xy, text, font, fill, anchor="mm", stroke=0):
    draw.text(xy, text, font=font, fill=fill, anchor=anchor, align="center",
              stroke_width=stroke, stroke_fill=fill)


def round_rect(draw, box, radius=16, fill="#FFFFFF", outline="#CDD7DF", width=2):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def build_preview():
    im = Image.open(TEMPLATE).convert("RGBA")
    draw = ImageDraw.Draw(im)
    # 본문 패널을 덮어 기존 예시 그림을 제거한다.
    round_rect(draw, (53, 157, 1560, 802), 19, "#F5F7F9", "#CFD6DC", 2)

    # PPT 레이아웃과 동일한 5단 흐름을 보여주는 고해상도 미리보기.
    cols = [(84, 332, "1", "시나리오·환경"), (362, 694, "2", "정책 생성"),
            (748, 974, "3", "일반화 검증"), (1012, 1285, "4", "정책 증류"),
            (1325, 1535, "5", "최종 비교")]
    for x1, x2, n, title in cols:
        draw.ellipse((x1, 190, x1 + 32, 222), fill="#0A83BF")
        draw_center(draw, (x1 + 16, 206), n, pil_font(17), "white")
        draw.text((x1 + 42, 188), title, font=pil_font(25), fill="#173A55")

    # 1단계
    round_rect(draw, (80, 244, 333, 693), 18)
    round_rect(draw, (123, 280, 178, 349), 8, "#EAF4FA", "#0A83BF", 2)
    round_rect(draw, (138, 270, 193, 339), 8, "#DDF0F8", "#0A83BF", 2)
    round_rect(draw, (154, 260, 209, 329), 8, "white", "#0A83BF", 2)
    draw_center(draw, (207, 370), "전국 시나리오 생성", pil_font(24), "#173A55")
    draw_center(draw, (207, 405), "시군구 · 병원 · AMB/UAV", pil_font(16), "#71808E")
    draw.line((207, 432, 207, 470), fill="#16A4D8", width=5)
    draw.polygon([(197, 464), (217, 464), (207, 479)], fill="#16A4D8")
    round_rect(draw, (163, 491, 251, 549), 10, "white", "#0A83BF", 2)
    draw.line((176, 521, 191, 521, 199, 507, 211, 537, 222, 519, 238, 519), fill="#17A398", width=3)
    draw_center(draw, (207, 577), "사건기반 시뮬레이션", pil_font(24), "#173A55")
    draw_center(draw, (207, 620), "동일 환경에서 정책 실행\n환자 이송–병원 치료 반영", pil_font(16), "#71808E")

    # 2단계
    round_rect(draw, (376, 250, 695, 402), 18, "#F7FAFC", "#BFCBD4", 2)
    round_rect(draw, (403, 286, 455, 350), 8, "white", "#71808E", 2)
    for yy in [301, 318, 335]:
        draw.line((413, yy, 420, yy + 6, 430, yy - 5), fill="#17A398", width=2)
        draw.line((437, yy + 2, 447, yy + 2), fill="#71808E", width=2)
    draw.text((475, 271), "64개 휴리스틱", font=pil_font(24), fill="#173A55")
    draw.text((475, 308), "규칙별 기준 성능 생성\n→ 비교군 확보", font=pil_font(16), fill="#71808E", spacing=5)
    round_rect(draw, (376, 426, 695, 693), 18, "#EEF8FC", "#7EC4E2", 3)
    draw.text((400, 448), "전국 RL 학습", font=pil_font(25), fill="#173A55")
    draw.text((400, 490), "시군구 다지역\nMaskable PPO", font=pil_font(16), fill="#71808E", spacing=5)
    agent = Image.open(AGENT).convert("RGBA")
    agent.thumbnail((128, 130), Image.Resampling.LANCZOS)
    im.alpha_composite(agent, (537, 454))
    round_rect(draw, (402, 607, 670, 659), 13, "#DDF2FA", "#DDF2FA", 1)
    draw_center(draw, (536, 633), "상태·보상  ↔  마스킹된 행동", pil_font(16), "#0A83BF")

    # 3단계
    round_rect(draw, (763, 266, 975, 669), 18)
    for d, c in [(66, "#0A83BF"), (44, "#16A4D8"), (22, "#F09A3E")]:
        draw.ellipse((869 - d // 2, 325 - d // 2, 869 + d // 2, 325 + d // 2), outline=c, width=3)
    draw_center(draw, (869, 393), "Holdout 평가", pil_font(25), "#173A55")
    draw_center(draw, (869, 440), "학습하지 않은\n지역·좌표", pil_font(17), "#71808E")
    round_rect(draw, (789, 500, 949, 552), 13, "#EEF3F6", "#EEF3F6", 1)
    draw_center(draw, (869, 526), "휴리스틱  ↔  RL", pil_font(17), "#0A83BF")
    draw_center(draw, (869, 604), "지역별 성능 지도\n동일 seed 통계 비교", pil_font(16), "#71808E")

    # 4단계
    round_rect(draw, (1024, 247, 1285, 693), 18)
    round_rect(draw, (1056, 283, 1107, 341), 7, "white", "#0A83BF", 2)
    for yy in [298, 312, 326]:
        draw.line((1066, yy, 1097, yy), fill="#71808E", width=2)
    draw.text((1120, 283), "RL 행동 로그", font=pil_font(21), fill="#173A55")
    draw.text((1120, 314), "상태–행동 수집", font=pil_font(15), fill="#71808E")
    # 트리
    nodes = [(1155, 412), (1128, 444), (1182, 444), (1113, 478), (1142, 478), (1168, 478), (1197, 478)]
    for a, b in [(0, 1), (0, 2), (1, 3), (1, 4), (2, 5), (2, 6)]:
        draw.line((*nodes[a], *nodes[b]), fill="#17A398", width=3)
    for i, (xx, yy) in enumerate(nodes):
        draw.ellipse((xx - 6, yy - 6, xx + 6, yy + 6), fill="#F09A3E" if i == 0 else "#17A398")
    draw_center(draw, (1155, 515), "VIPER 정책 증류", pil_font(22), "#173A55")
    draw_center(draw, (1155, 548), "RL → 의사결정트리", pil_font(16), "#71808E")
    round_rect(draw, (1060, 607, 1253, 660), 13, "#E7F6F3", "#B5DDD7", 2)
    draw_center(draw, (1156, 634), "Tree 재시뮬레이션", pil_font(17), "#17A398")

    # 5단계
    round_rect(draw, (1334, 254, 1535, 681), 18)
    for i, (w, c) in enumerate([(72, "#71808E"), (94, "#0A83BF"), (82, "#17A398")]):
        round_rect(draw, (1390, 296 + i * 30, 1390 + w, 309 + i * 30), 6, c, c, 1)
    draw_center(draw, (1435, 407), "3자 최종 비교", pil_font(23), "#173A55")
    for i, (label, fill, color) in enumerate([("휴리스틱", "#EEF1F3", "#71808E"), ("RL 정책", "#E5F3FA", "#0A83BF"), ("VIPER Tree", "#E4F5F2", "#17A398")]):
        round_rect(draw, (1373, 452 + i * 62, 1518, 496 + i * 62), 12, fill, fill, 1)
        draw_center(draw, (1445, 474 + i * 62), label, pil_font(17), color)
    draw_center(draw, (1435, 650), "생존 성과 · 일반화\n· 해석성", pil_font(15), "#71808E")

    # 시뮬레이터 출력은 휴리스틱/RL로 분기하고, 두 정책은 Holdout에서 다시 합류한다.
    draw.line((333, 566, 351, 566), fill="#16A4D8", width=6)
    draw.line((351, 326, 351, 566), fill="#16A4D8", width=6)
    for yy in [326, 560]:
        draw.line((351, yy, 369, yy), fill="#16A4D8", width=6)
        draw.polygon([(363, yy - 9), (380, yy), (363, yy + 9)], fill="#16A4D8")
    draw.line((695, 326, 716, 326), fill="#16A4D8", width=6)
    draw.line((695, 560, 716, 560), fill="#16A4D8", width=6)
    draw.line((716, 326, 716, 560), fill="#16A4D8", width=6)
    draw.line((716, 440, 748, 440), fill="#16A4D8", width=6)
    draw.polygon([(742, 431), (759, 440), (742, 449)], fill="#16A4D8")

    # Holdout → 로그/증류 → 재실행 → 최종 비교의 순서를 유지한다.
    for x in [980, 1290]:
        draw.line((x, 440, x + 28, 440), fill="#16A4D8", width=6)
        draw.polygon([(x + 22, 431), (x + 39, 440), (x + 22, 449)], fill="#16A4D8")
    for y1, y2 in [(350, 389), (563, 596)]:
        draw.line((1155, y1, 1155, y2), fill="#16A4D8", width=4)
        draw.polygon([(1147, y2 - 8), (1163, y2 - 8), (1155, y2 + 4)], fill="#16A4D8")

    # 하단 핵심 문장
    round_rect(draw, (100, 722, 1514, 774), 14, "#E8F3F8", "#C7E2EF", 1)
    draw_center(draw, (807, 748), "전국 시나리오에서 학습 → 미학습 지역에서 검증 → 해석 가능한 정책으로 환류", pil_font(21), "#173A55")

    im.convert("RGB").save(PNG_OUT, quality=95)


if __name__ == "__main__":
    OUT.mkdir(parents=True, exist_ok=True)
    build_pptx()
    build_preview()
    print(PPTX_OUT)
    print(PNG_OUT)
